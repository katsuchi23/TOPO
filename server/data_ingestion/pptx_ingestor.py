from pptx import Presentation
import re
import json

class PPTXIngestor:
    def __init__(self, file_path):
        self.file_path = file_path
        
    def ingest(self):
        prs = Presentation(self.file_path)
        result = {}
        
        for slide in prs.slides:
            for shape in slide.shapes:
                # Parse tables (Quarterly Metrics)
                if shape.has_table:
                    table = shape.table
                    rows = list(table.rows)
                    quarterly_data = []
                    for row in rows[1:5]:  # Process Q1-Q4 rows
                        cells = [cell.text.strip() for cell in row.cells]
                        quarterly_data.append({
                            "quarter": cells[0],
                            "revenue": float(cells[1].replace('$', '').replace(',', '')),
                            "memberships_sold": int(cells[2].replace(',', '')),
                            "avg_duration_minutes": int(cells[3])
                        })
                    result["quarterly_metrics"] = quarterly_data
                
                # Parse text boxes
                if hasattr(shape, "text"):
                    text = shape.text
                    lines = text.split('\n')
                    for line in lines:
                        line = line.strip()
                        if not line:
                            continue
                        
                        # Total Revenue
                        if line.startswith('Total Revenue:'):
                            result['total_revenue'] = float(re.search(r'\$([\d,]+)', line).group(1).replace(',', ''))
                        
                        # Total Memberships
                        elif line.startswith('Total Memberships Sold:'):
                            result['total_memberships'] = int(re.search(r'\d+', line).group())
                        
                        # Top Location
                        elif line.startswith('Top Location:'):
                            result['top_location'] = line.split(': ')[1].strip()
                        
                        # Revenue Distribution
                        elif line == 'Revenue Distribution:':
                            result['revenue_distribution'] = {}
                        elif 'revenue_distribution' in result and ': ' in line:
                            key, value = line.split(': ', 1)
                            key_clean = key.strip().lower().replace(' ', '_')
                            result['revenue_distribution'][key_clean] = float(value.replace('%', ''))

        # Calculate and enforce correct totals from quarterly data
        if 'quarterly_metrics' in result:
            # Calculate correct values
            correct_memberships = sum(q['memberships_sold'] for q in result['quarterly_metrics'])
            correct_revenue = sum(q['revenue'] for q in result['quarterly_metrics'])
            
            # Check and update if needed
            if result.get('total_memberships') != correct_memberships:
                print(f"Correcting total_memberships: {result.get('total_memberships')} → {correct_memberships}")
                result['total_memberships'] = correct_memberships
                
            if result.get('total_revenue') != correct_revenue:
                print(f"Correcting total_revenue: {result.get('total_revenue')} → {correct_revenue}")
                result['total_revenue'] = correct_revenue
        
        return result